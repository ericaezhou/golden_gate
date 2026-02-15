#!/usr/bin/env python3
"""Generate demo data files for Golden Gate / Bridge AI.

Creates 8 realistic work artifacts for the Alice Chen risk analyst scenario.
All files are generated programmatically and placed in data/.

Usage:
    python scripts/generate_demo_data.py
    python scripts/generate_demo_data.py --output-dir /custom/path

Gaps embedded:
    1. Macro Overlay Black Box - undocumented manual overlay criteria
    2. Missing Escalation Playbook - undefined/inconsistent escalation thresholds
    3. Legacy Shortcut vs Policy Method - 4q vs 12q, undocumented switching criteria
"""

from __future__ import annotations

import argparse
import os
import shutil
import sqlite3
from pathlib import Path

import nbformat
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.workbook.defined_name import DefinedName
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt


ROOT = Path(__file__).resolve().parent.parent
DEFAULT_OUTPUT = ROOT / "data"


# ---------------------------------------------------------------------------
# 1. loss_forecast_model.py
# ---------------------------------------------------------------------------
def generate_python_model(output_dir: Path) -> None:
    content = '''\
"""
Credit Loss Forecasting Model — Q3 2024
Author: Alice Chen, Risk Analyst
Last Updated: 2024-09-15

Calculates quarterly credit loss forecasts using historical defaults
and macroeconomic adjustments.

References:
    - Data source: portfolio_risk.db
    - SQL queries: risk_queries.sql
    - Policy doc:  model_methodology.docx
    - Output:      Q3_2024_forecast.xlsx
"""

import sqlite3
import pandas as pd
import numpy as np
from datetime import datetime
from typing import Dict, Optional, Tuple

# --- Configuration ---
DB_PATH = "portfolio_risk.db"
FORECAST_OUTPUT = "Q3_2024_forecast.xlsx"

# TODO: Policy (model_methodology.docx) specifies 12 quarters.
#       Using 4 here because Alice says the difference is < 0.3%
#       when macro conditions are stable.
LOOKBACK_QUARTERS = 4

# Segment definitions
SEGMENTS = ["prime", "near_prime", "subprime", "deep_subprime"]

# Macro overlay multiplier — applied when GDP drops below trigger
# See policy_snapshot table in portfolio_risk.db for approved value
MACRO_OVERLAY_MULTIPLIER = None  # TODO: Ask Alice — she applies this manually

# Escalation thresholds (dollar impact)
ESCALATION_THRESHOLDS = {
    "routine": 0,               # Analyst sign-off only
    "risk_committee": 2_000_000, # Weekly sync review
    "cfo_approval": 5_000_000,   # CFO email approval (3-5 day lead time)
    "board_notification": None,  # Verbal guidance from CFO (~$10M) — not in any policy
}

# Segment-level loss thresholds for triggering an overlay
SEGMENT_LOSS_THRESHOLDS = {
    "prime": 0.03,
    "near_prime": 0.07,
    "subprime": None,        # Alice adjusts this quarterly — no fixed value
    "deep_subprime": None,   # Not defined; Alice decides at runtime
}


def load_historical_defaults(db_path: str = DB_PATH) -> pd.DataFrame:
    """Load historical default rates from portfolio_risk.db."""
    conn = sqlite3.connect(db_path)
    df = pd.read_sql("SELECT * FROM historical_defaults ORDER BY quarter DESC", conn)
    conn.close()
    return df


def compute_baseline_rate(
    defaults_df: pd.DataFrame,
    n_quarters: int = LOOKBACK_QUARTERS,
) -> float:
    """Compute rolling average default rate.

    NOTE: Policy specifies 12-quarter window (see risk_queries.sql query A).
    We use 4-quarter here for speed. Alice says it\'s fine when GDP > 1.5%.
    """
    recent = defaults_df.head(n_quarters)
    return recent["default_rate"].mean()


def apply_macro_overlay(
    base_rate: float,
    macro_data: Optional[Dict] = None,
) -> Tuple[float, float]:
    """Apply macroeconomic adjustment to base rate.

    Returns (adjusted_rate, overlay_amount).

    WARNING: The overlay sizing logic below is a rough sketch.
    The actual criteria Alice uses are NOT fully captured here.
    """
    if macro_data is None:
        return base_rate, 0.0

    overlay = 0.0

    # TODO: Manual overlay criteria not documented.
    #       Alice applies +1% to +5% based on her judgment.
    #       Rough heuristic from her notes:
    #         - 10% delinquency increase → +1% overlay
    #         - Cap at 5% per segment
    #       But she also considers cohort variance, new product buffers,
    #       and "gut feel" from 4 years of experience.
    # FIXME: This logic MUST be formalized before Alice leaves.

    if macro_data.get("gdp_growth", 0) < 1.5:
        overlay += 0.02  # Recession buffer
    if macro_data.get("unemployment_delta", 0) > 0.5:
        overlay += macro_data["unemployment_delta"] * 0.01

    return base_rate + overlay, overlay


def should_escalate(impact_amount: float) -> str:
    """Determine escalation level based on dollar impact.

    Returns one of: \'routine\', \'risk_committee\', \'cfo\', \'board\'.

    NOTE: Board notification threshold was never formally approved.
    Alice operates on verbal ~$10M guidance from CFO.
    """
    # TODO: Board threshold is undefined — Alice handles this case manually
    if impact_amount >= 10_000_000:
        return "board"  # Alice\'s verbal rule, not in policy
    for level in ["cfo_approval", "risk_committee", "routine"]:
        threshold = ESCALATION_THRESHOLDS.get(level, 0) or 0
        if impact_amount >= threshold:
            return level.replace("_approval", "")
    return "routine"


def generate_forecast(db_path: str = DB_PATH) -> pd.DataFrame:
    """Generate full quarterly forecast.

    Output is saved to Q3_2024_forecast.xlsx after manual overlay review.
    """
    defaults_df = load_historical_defaults(db_path)
    baseline = compute_baseline_rate(defaults_df)

    results = []
    for segment in SEGMENTS:
        threshold = SEGMENT_LOSS_THRESHOLDS.get(segment)
        results.append({
            "segment": segment,
            "baseline_rate": round(baseline, 4),
            "macro_overlay": 0.0,        # Filled in manually by Alice
            "final_rate": round(baseline, 4),  # Updated after overlay
            "threshold": threshold,       # May be None — that\'s a problem
            "rationale": "",              # Alice fills this in (but rarely does)
        })

    return pd.DataFrame(results)


if __name__ == "__main__":
    forecast = generate_forecast()
    forecast.to_excel(FORECAST_OUTPUT, index=False)
    print(forecast)
'''
    (output_dir / "loss_forecast_model.py").write_text(content)


# ---------------------------------------------------------------------------
# 2. Q3_2024_forecast.xlsx
# ---------------------------------------------------------------------------
def generate_excel_forecast(output_dir: Path) -> None:
    wb = Workbook()

    # --- Sheet 1: Forecast_Summary ---
    ws = wb.active
    ws.title = "Forecast_Summary"
    headers = ["Segment", "Model Rate", "Macro Overlay", "Final Rate",
               "Reserve Impact ($M)", "Rationale"]
    ws.append(headers)

    data = [
        ("Prime",          0.023, 0.000, None, None, "No adjustment needed"),
        ("Near Prime",     0.051, 0.012, None, None, ""),
        ("Subprime",       0.124, 0.031, None, None, ""),
        ("Deep Subprime",  0.221, 0.045, None, None, ""),
    ]
    exposures = [145_000_000, 82_000_000, 38_000_000, 12_000_000]

    for i, (seg, model, overlay, _, _, rationale) in enumerate(data, start=2):
        ws.cell(row=i, column=1, value=seg)
        ws.cell(row=i, column=2, value=model)
        ws.cell(row=i, column=3, value=overlay)
        ws.cell(row=i, column=4).value = f"=B{i}+C{i}"  # formula
        ws.cell(row=i, column=5).value = f"=D{i}*{exposures[i-2]}"  # formula
        ws.cell(row=i, column=6, value=rationale)

    # Total row
    ws.cell(row=6, column=1, value="TOTAL")
    ws.cell(row=6, column=5).value = "=SUM(E2:E5)"

    # Comments
    ws["C3"].comment = Comment("Based on early delinquency signals - AC", "Alice Chen")
    ws["C4"].comment = Comment("Elevated 30-day DQ + macro concerns - AC", "Alice Chen")
    ws["C5"].comment = Comment("New product buffer + macro environment - AC", "Alice Chen")
    ws["F3"].comment = Comment("TODO: document rationale before quarter-end", "Alice Chen")

    # Named range
    dn = DefinedName("OVERLAY_RANGE", attr_text="Forecast_Summary!$C$2:$C$5")
    wb.defined_names.add(dn)

    # --- Sheet 2: Cohort_Tracking ---
    ws2 = wb.create_sheet("Cohort_Tracking")
    ws2.append(["Quarter", "Vintage", "Segment", "Expected", "Actual", "Variance", "Status"])
    cohorts = [
        ("Q1-24", "Jan", "Near Prime",  0.042, 0.058, None, "REVIEW"),
        ("Q1-24", "Feb", "Near Prime",  0.041, 0.054, None, "REVIEW"),
        ("Q1-24", "Mar", "Near Prime",  0.040, 0.043, None, "OK"),
        ("Q4-23", "Oct", "Subprime",    0.105, 0.132, None, "REVIEW"),
        ("Q4-23", "Nov", "Subprime",    0.108, 0.119, None, "MONITOR"),
        ("Q3-23", "Jul", "Deep Sub",    0.195, 0.261, None, "REVIEW"),
        ("Q3-23", "Aug", "Deep Sub",    0.201, 0.247, None, "REVIEW"),
    ]
    for i, (q, v, seg, exp, act, _, status) in enumerate(cohorts, start=2):
        ws2.cell(row=i, column=1, value=q)
        ws2.cell(row=i, column=2, value=v)
        ws2.cell(row=i, column=3, value=seg)
        ws2.cell(row=i, column=4, value=exp)
        ws2.cell(row=i, column=5, value=act)
        ws2.cell(row=i, column=6).value = f"=(E{i}-D{i})/D{i}"
        ws2.cell(row=i, column=7, value=status)

    ws2["E2"].comment = Comment(
        "Two months >25% variance = overlay trigger per Alice's rule", "Alice Chen"
    )

    # --- Sheet 3: Adjustment_Log ---
    ws3 = wb.create_sheet("Adjustment_Log")
    ws3.append(["Date", "Segment", "Overlay Applied", "Trigger", "Approved By", "Removal Date"])
    log_data = [
        ("2024-06-15", "Near Prime",     "+1.2%", "Q1 cohort variance >25%",    "", ""),
        ("2024-06-15", "Subprime",       "+3.1%", "DQ trend + cohort miss",     "", ""),
        ("2024-03-01", "Deep Subprime",  "+4.5%", "New product buffer",         "", ""),
        ("2023-12-01", "Subprime",       "+2.8%", "Macro deterioration",        "", ""),
    ]
    for row in log_data:
        ws3.append(list(row))

    ws3["E2"].comment = Comment(
        "Approval column is frequently left blank — governance gap", "System"
    )

    # --- Sheet 4: Alice_Notes (HIDDEN) ---
    ws4 = wb.create_sheet("Alice_Notes")
    ws4.sheet_state = "hidden"
    ws4.append(["Parameter", "Value"])
    notes = [
        ("Real subprime threshold",  "0.18 (not in any config file)"),
        ("Real deep_sub threshold",  "0.28 (I set this based on vintage performance)"),
        ("Overlay cap",              "5% max per segment"),
        ("Board escalation threshold", "$10M — verbal from CFO, not in policy"),
        ("When legacy shortcut is OK", "GDP growth > 1.5% and no Fed rate changes"),
        ("Overlay removal rule",     "3 consecutive months of <10% variance"),
    ]
    for row in notes:
        ws4.append(list(row))
    ws4["A1"].comment = Comment("Personal reference — not for distribution", "Alice Chen")

    wb.save(output_dir / "Q3_2024_forecast.xlsx")


# ---------------------------------------------------------------------------
# 3. portfolio_risk.db
# ---------------------------------------------------------------------------
def generate_sqlite_db(output_dir: Path) -> None:
    db_path = output_dir / "portfolio_risk.db"
    if db_path.exists():
        db_path.unlink()

    conn = sqlite3.connect(str(db_path))
    c = conn.cursor()

    # Table: historical_defaults (16 rows = 4 years, Q1-2021 to Q4-2024)
    c.execute("""
        CREATE TABLE historical_defaults (
            quarter TEXT PRIMARY KEY,
            default_rate REAL NOT NULL
        )
    """)
    defaults = [
        ("Q1-2021", 0.032), ("Q2-2021", 0.034), ("Q3-2021", 0.031), ("Q4-2021", 0.035),
        ("Q1-2022", 0.038), ("Q2-2022", 0.041), ("Q3-2022", 0.039), ("Q4-2022", 0.043),
        ("Q1-2023", 0.048), ("Q2-2023", 0.050), ("Q3-2023", 0.052), ("Q4-2023", 0.055),
        ("Q1-2024", 0.056), ("Q2-2024", 0.056), ("Q3-2024", 0.058), ("Q4-2024", 0.060),
    ]
    c.executemany("INSERT INTO historical_defaults VALUES (?, ?)", defaults)

    # Table: loan_cohorts (4 segments)
    c.execute("""
        CREATE TABLE loan_cohorts (
            segment TEXT PRIMARY KEY,
            exposure_ead REAL NOT NULL,
            lgd REAL NOT NULL,
            avg_fico INTEGER,
            avg_dti REAL
        )
    """)
    cohorts = [
        ("prime",          145_000_000, 0.38, 744, 0.22),
        ("near_prime",      82_000_000, 0.47, 692, 0.28),
        ("subprime",        38_000_000, 0.58, 622, 0.34),
        ("deep_subprime",   12_000_000, 0.68, 565, 0.41),
    ]
    c.executemany("INSERT INTO loan_cohorts VALUES (?, ?, ?, ?, ?)", cohorts)

    # Table: macro_quarterly (4 recent quarters)
    c.execute("""
        CREATE TABLE macro_quarterly (
            quarter TEXT PRIMARY KEY,
            gdp_growth REAL,
            unemployment_rate REAL,
            fed_funds_rate REAL,
            consumer_sentiment REAL
        )
    """)
    macro = [
        ("Q1-2024", 1.6, 3.9, 5.25, 67.4),
        ("Q2-2024", 1.4, 4.0, 5.25, 65.2),
        ("Q3-2024", 1.1, 4.2, 5.00, 63.8),
        ("Q4-2024", 0.8, 4.5, 4.75, 61.1),
    ]
    c.executemany("INSERT INTO macro_quarterly VALUES (?, ?, ?, ?, ?)", macro)

    # Table: escalation_thresholds — board row has NULLs (Gap 2)
    c.execute("""
        CREATE TABLE escalation_thresholds (
            level TEXT PRIMARY KEY,
            dollar_threshold REAL,
            approval_required TEXT,
            last_updated TEXT
        )
    """)
    thresholds = [
        ("routine",        0,         "analyst",        "2024-01-15"),
        ("risk_committee", 2_000_000, "risk_committee", "2024-01-15"),
        ("cfo",            5_000_000, "cfo",            "2024-01-15"),
        ("board",          None,       None,             None),  # NEVER FORMALLY SET
    ]
    c.executemany("INSERT INTO escalation_thresholds VALUES (?, ?, ?, ?)", thresholds)

    # Table: policy_snapshot — says 12-quarter window (Gap 3)
    c.execute("""
        CREATE TABLE policy_snapshot (
            as_of_quarter TEXT PRIMARY KEY,
            lookback_window INTEGER,
            overlay_cap REAL,
            methodology TEXT
        )
    """)
    c.execute(
        "INSERT INTO policy_snapshot VALUES (?, ?, ?, ?)",
        ("Q3-2024", 12, 0.05, "12-quarter rolling average per CR-POL-2024-003"),
    )

    # Table: capital_position (4 quarters)
    c.execute("""
        CREATE TABLE capital_position (
            quarter TEXT PRIMARY KEY,
            total_capital REAL,
            risk_weighted_assets REAL,
            capital_ratio REAL
        )
    """)
    capital = [
        ("Q1-2024", 2_750_000_000, 21_200_000_000, 0.1297),
        ("Q2-2024", 2_758_000_000, 21_270_000_000, 0.1297),
        ("Q3-2024", 2_787_000_000, 21_320_000_000, 0.1307),
        ("Q4-2024", 2_763_000_000, 21_430_000_000, 0.1289),
    ]
    c.executemany("INSERT INTO capital_position VALUES (?, ?, ?, ?)", capital)

    conn.commit()
    conn.close()


# ---------------------------------------------------------------------------
# 4. risk_queries.sql
# ---------------------------------------------------------------------------
def generate_sql_queries(output_dir: Path) -> None:
    content = """\
-- risk_queries.sql
-- Q3 2024 Credit Loss Forecast Workstream
-- Author: Alice Chen
--
-- NOTE: Two query variants are included for the default rate calculation:
--   (A) POLICY_COMPLIANT — 12-quarter rolling average (per CR-POL-2024-003)
--   (B) LEGACY_SHORTCUT  — 4-quarter rolling average
--       Commonly used by Alice historically. Faster, but may violate policy.
--       See Alice for when to use (A) vs (B).
--
-- Data source: portfolio_risk.db
-- See also: loss_forecast_model.py, model_methodology.docx

-----------------------------------------------------------------------
-- POLICY SNAPSHOT (verify current parameters)
-----------------------------------------------------------------------
SELECT *
FROM policy_snapshot
WHERE as_of_quarter = 'Q3-2024';

-----------------------------------------------------------------------
-- GDP TRIGGER CHECK
-----------------------------------------------------------------------
SELECT quarter, gdp_growth
FROM macro_quarterly
WHERE quarter = 'Q3-2024';

-----------------------------------------------------------------------
-- (A) POLICY_COMPLIANT: 12-quarter rolling average default rate
-----------------------------------------------------------------------
WITH recent AS (
    SELECT quarter, default_rate
    FROM historical_defaults
    ORDER BY quarter DESC
    LIMIT 12
)
SELECT AVG(default_rate) AS avg_default_rate_12q
FROM recent;

-----------------------------------------------------------------------
-- (B) LEGACY_SHORTCUT: 4-quarter rolling average (NOT policy-compliant)
--     Alice uses this when GDP > 1.5% and macro is "stable."
--     Difference vs 12q is typically < 0.3%, but can reach 1-2%
--     during regime changes.
-----------------------------------------------------------------------
WITH recent AS (
    SELECT quarter, default_rate
    FROM historical_defaults
    ORDER BY quarter DESC
    LIMIT 4
)
SELECT AVG(default_rate) AS avg_default_rate_4q
FROM recent;

-----------------------------------------------------------------------
-- LATEST CAPITAL POSITION
-----------------------------------------------------------------------
SELECT *
FROM capital_position
ORDER BY quarter DESC
LIMIT 1;

-----------------------------------------------------------------------
-- COHORT EXPOSURES (EAD) and LGD by segment
-----------------------------------------------------------------------
SELECT segment, exposure_ead, lgd, avg_fico
FROM loan_cohorts;

-----------------------------------------------------------------------
-- ESCALATION THRESHOLDS
-- NOTE: board row has NULL thresholds — see loss_forecast_model.py
-----------------------------------------------------------------------
SELECT *
FROM escalation_thresholds
ORDER BY dollar_threshold;
"""
    (output_dir / "risk_queries.sql").write_text(content)


# ---------------------------------------------------------------------------
# 5. board_risk_presentation.pptx
# ---------------------------------------------------------------------------
def generate_pptx_deck(output_dir: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Q3 2024 Credit Risk Update"
    slide.placeholders[1].text = (
        "Board Risk Committee\nPrepared by Alice Chen, Risk Analyst"
    )
    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "Standard quarterly update. Model outputs from loss_forecast_model.py, "
        "adjustments in Q3_2024_forecast.xlsx. Methodology in model_methodology.docx."
    )

    # Slide 2: Forecast Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Portfolio Loss Forecast"
    # Add table
    rows, cols = 6, 5
    tbl_shape = slide2.shapes.add_table(rows, cols, PptxInches(1), PptxInches(2),
                                         PptxInches(11), PptxInches(3.5))
    tbl = tbl_shape.table
    headers = ["Segment", "Model Rate", "Adjusted Rate", "Reserve Impact ($M)", "Overlay"]
    for j, h in enumerate(headers):
        tbl.cell(0, j).text = h
    data = [
        ("Prime",         "2.3%",  "2.3%",  "$3.3M",  "0.0%"),
        ("Near Prime",    "5.1%",  "6.3%",  "$5.2M",  "+1.2%"),
        ("Subprime",      "12.4%", "15.5%", "$5.9M",  "+3.1%"),
        ("Deep Subprime", "22.1%", "26.6%", "$3.2M",  "+4.5%"),
        ("TOTAL",         "",      "",      "$17.6M", ""),
    ]
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            tbl.cell(i, j).text = val

    notes2 = slide2.notes_slide
    notes2.notes_text_frame.text = (
        "CFO asked about adjustment methodology last quarter. Need to document "
        "the overlay criteria. See model_methodology.docx Sections 4-5 (currently incomplete). "
        "Source: Q3_2024_forecast.xlsx"
    )

    # Slide 3: Macro Environment
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Macroeconomic Environment & Adjustments"
    body = slide3.placeholders[1]
    tf = body.text_frame
    tf.text = "Macroeconomic factors considered:"
    for bullet in [
        "GDP growth trending down (1.6% → 0.8% over 2024)",
        "Unemployment rising (3.9% → 4.5%)",
        "Fed rate easing cycle beginning (5.25% → 4.75%)",
        "Consumer sentiment declining (67.4 → 61.1)",
    ]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "\nOverlays reflect analyst judgment based on these indicators."

    notes3 = slide3.notes_slide
    notes3.notes_text_frame.text = (
        "Don't mention the 4-quarter shortcut to the board. Policy says 12 quarters "
        "but the difference is minimal in current conditions. If asked about methodology, "
        "refer to model_methodology.docx. The board deck numbers are based on the "
        "4-quarter baseline from loss_forecast_model.py."
    )

    # Slide 4: Escalation & Governance
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Escalation & Governance"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Current escalation framework:"
    for bullet in [
        "< $2M impact: Analyst sign-off (routine)",
        "$2M - $5M impact: Risk Committee review",
        "> $5M impact: CFO approval required",
        "Board notification threshold: Under review",
    ]:
        p = tf4.add_paragraph()
        p.text = bullet
        p.level = 1
    p3 = tf4.add_paragraph()
    p3.text = "\nAll overlays reviewed at weekly Risk Committee sync."

    notes4 = slide4.notes_slide
    notes4.notes_text_frame.text = (
        "Board notification threshold still not formally approved. Operating on "
        "verbal guidance from CFO (~$10M). Should formalize before Alice's transition. "
        "See escalation_thresholds table in portfolio_risk.db — the board row is NULL."
    )

    # Slide 5: Next Steps
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Action Items"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Key action items:"
    for bullet in [
        "Model recalibration scheduled for Q1 2025",
        "Document overlay decision criteria (model_methodology.docx Sections 4-5)",
        "Formalize board notification threshold in policy",
        "Transition planning: train Marcus Park as backup analyst",
        "Reconcile legacy vs. policy-compliant parameter differences",
    ]:
        p = tf5.add_paragraph()
        p.text = bullet
        p.level = 1

    notes5 = slide5.notes_slide
    notes5.notes_text_frame.text = (
        "The transition planning bullet was added at CFO's request. Priority is getting "
        "overlay methodology documented before Alice's last day. Marcus has been shadowing "
        "for 2 months but hasn't run overlays independently. "
        "See stress_testing.ipynb for scenario analysis approach."
    )

    prs.save(output_dir / "board_risk_presentation.pptx")


# ---------------------------------------------------------------------------
# 6. model_methodology.docx
# ---------------------------------------------------------------------------
def generate_docx_methodology(output_dir: Path) -> None:
    doc = Document()

    # Title
    doc.add_heading("Credit Loss Forecasting Model — Methodology Document", level=0)

    # Metadata table
    meta = doc.add_table(rows=5, cols=2)
    meta.style = "Table Grid"
    fields = [
        ("Document ID", "CR-METH-2024-001"),
        ("Version", "1.3 (DRAFT)"),
        ("Author", "Alice Chen"),
        ("Last Review", "July 2024"),
        ("Status", "INCOMPLETE — Sections 4 and 5 require documentation"),
    ]
    for i, (k, v) in enumerate(fields):
        meta.cell(i, 0).text = k
        meta.cell(i, 1).text = v

    doc.add_paragraph()

    # Section 1: Overview
    doc.add_heading("1. Overview", level=1)
    doc.add_paragraph(
        "This document describes the methodology for the quarterly credit loss "
        "forecasting model (loss_forecast_model.py). The model generates baseline "
        "loss rate projections which are then subject to analyst judgment overlays."
    )
    doc.add_paragraph(
        "Source data is maintained in portfolio_risk.db. Final forecasts are "
        "published in Q3_2024_forecast.xlsx. Reference SQL queries are in "
        "risk_queries.sql."
    )

    # Section 2: Baseline Calculation
    doc.add_heading("2. Baseline Calculation", level=1)
    doc.add_paragraph(
        "The baseline loss rate is computed as a rolling average of historical "
        "default rates from the historical_defaults table in portfolio_risk.db."
    )
    p = doc.add_paragraph()
    run = p.add_run(
        "Per policy CR-POL-2024-003, the rolling window shall be 12 quarters."
    )
    run.bold = True

    doc.add_paragraph(
        "See risk_queries.sql for the reference implementation (Query A)."
    )
    doc.add_paragraph(
        "NOTE: An alternative 4-quarter window calculation exists in the codebase "
        "(risk_queries.sql Query B, loss_forecast_model.py). This was used historically "
        "for preliminary estimates. The 12-quarter window is the policy-compliant method "
        "and should be used for all official forecasts."
    )

    # Section 3: Macroeconomic Adjustments
    doc.add_heading("3. Macroeconomic Adjustments", level=1)
    doc.add_paragraph(
        "The model incorporates macroeconomic indicators to adjust forecasts "
        "when conditions deviate from baseline assumptions."
    )

    tbl = doc.add_table(rows=5, cols=3)
    tbl.style = "Table Grid"
    tbl.cell(0, 0).text = "Indicator"
    tbl.cell(0, 1).text = "Source"
    tbl.cell(0, 2).text = "Update Frequency"
    indicators = [
        ("GDP growth rate", "BEA", "Quarterly"),
        ("Unemployment rate", "BLS", "Monthly"),
        ("Federal funds rate", "Federal Reserve", "As announced"),
        ("Consumer sentiment", "U. of Michigan", "Monthly"),
    ]
    for i, (ind, src, freq) in enumerate(indicators, start=1):
        tbl.cell(i, 0).text = ind
        tbl.cell(i, 1).text = src
        tbl.cell(i, 2).text = freq

    # Section 4: Overlay Decision Framework (INCOMPLETE — Gap 1)
    doc.add_heading("4. Overlay Decision Framework", level=1)
    p4 = doc.add_paragraph()
    run4 = p4.add_run("[SECTION INCOMPLETE]")
    run4.bold = True

    doc.add_paragraph(
        "Manual overlays are applied when model output diverges from observed "
        "market conditions. For the complete overlay decision criteria, "
        "including trigger thresholds, sizing methodology, and removal criteria, "
        "see Alice Chen."
    )
    doc.add_paragraph("General guidance:")
    doc.add_paragraph("Overlays should be supported by quantitative evidence", style="List Bullet")
    doc.add_paragraph(
        "Maximum overlay per segment: see SEGMENT_LOSS_THRESHOLDS in loss_forecast_model.py",
        style="List Bullet",
    )
    doc.add_paragraph("Overlay documentation is required per escalation policy", style="List Bullet")
    doc.add_paragraph(
        "Detailed criteria for when to apply overlays, how to size them, and "
        "when to remove them are maintained by the primary risk analyst and have "
        "not been formally documented."
    )

    # Section 5: Escalation Thresholds (INCOMPLETE — Gap 2)
    doc.add_heading("5. Escalation Thresholds", level=1)
    p5 = doc.add_paragraph()
    run5 = p5.add_run("[SECTION INCOMPLETE]")
    run5.bold = True

    doc.add_paragraph(
        "Threshold configuration is in loss_forecast_model.py. Current thresholds:"
    )

    etbl = doc.add_table(rows=5, cols=3)
    etbl.style = "Table Grid"
    etbl.cell(0, 0).text = "Level"
    etbl.cell(0, 1).text = "Dollar Impact"
    etbl.cell(0, 2).text = "Approval Required"
    esc_data = [
        ("Routine", "< $2M", "Analyst sign-off"),
        ("Significant", "$2M – $5M", "Risk Committee"),
        ("Major", "> $5M", "CFO approval"),
        ("Board notification", "See Alice Chen", "TBD"),
    ]
    for i, (lvl, impact, approval) in enumerate(esc_data, start=1):
        etbl.cell(i, 0).text = lvl
        etbl.cell(i, 1).text = impact
        etbl.cell(i, 2).text = approval

    doc.add_paragraph(
        "Note: Segment-level loss rate thresholds that trigger escalation "
        "review are maintained dynamically by Alice Chen and are not currently "
        "hardcoded in the system."
    )

    # Section 6: Stress Testing
    doc.add_heading("6. Stress Testing", level=1)
    doc.add_paragraph(
        "Stress scenarios are run using stress_testing.ipynb. The notebook applies "
        "macro stress shocks to the baseline forecast and generates scenario-specific "
        "projections for Risk Committee review."
    )

    # Section 7: Review History
    doc.add_heading("7. Review History", level=1)
    rtbl = doc.add_table(rows=5, cols=3)
    rtbl.style = "Table Grid"
    rtbl.cell(0, 0).text = "Version"
    rtbl.cell(0, 1).text = "Date"
    rtbl.cell(0, 2).text = "Changes"
    revisions = [
        ("1.0", "2023-06-01", "Initial draft"),
        ("1.1", "2023-12-15", "Added macro adjustment section"),
        ("1.2", "2024-03-01", "Updated escalation thresholds"),
        ("1.3", "2024-07-15", "Marked Sections 4 & 5 incomplete"),
    ]
    for i, (ver, dt, changes) in enumerate(revisions, start=1):
        rtbl.cell(i, 0).text = ver
        rtbl.cell(i, 1).text = dt
        rtbl.cell(i, 2).text = changes

    doc.save(output_dir / "model_methodology.docx")


# ---------------------------------------------------------------------------
# 7. stress_testing.ipynb
# ---------------------------------------------------------------------------
def generate_notebook(output_dir: Path) -> None:
    nb = nbformat.v4.new_notebook()
    nb.metadata["kernelspec"] = {
        "display_name": "Python 3",
        "language": "python",
        "name": "python3",
    }

    # Cell 1: Title (markdown)
    nb.cells.append(nbformat.v4.new_markdown_cell(
        "# Stress Testing — Q3 2024\n\n"
        "Applies macro stress scenarios to baseline loss forecast.\n\n"
        "Data source: `portfolio_risk.db`\n"
        "Model reference: `loss_forecast_model.py`\n"
        "Policy reference: `model_methodology.docx`"
    ))

    # Cell 2: Imports (code)
    c2 = nbformat.v4.new_code_cell(
        "import sqlite3\n"
        "import pandas as pd\n"
        "import numpy as np\n\n"
        "DB_PATH = \"portfolio_risk.db\"\n"
        "conn = sqlite3.connect(DB_PATH)"
    )
    nb.cells.append(c2)

    # Cell 3: Load defaults + compute baseline (code)
    c3 = nbformat.v4.new_code_cell(
        "# Load historical defaults\n"
        "defaults = pd.read_sql(\n"
        "    \"SELECT * FROM historical_defaults ORDER BY quarter DESC\", conn\n"
        ")\n\n"
        "# Use 4-quarter average for baseline (faster iteration during stress runs)\n"
        "baseline_rate = defaults.head(4)[\"default_rate\"].mean()\n"
        "print(f\"Baseline rate (4Q avg): {baseline_rate:.4f}\")"
    )
    c3.outputs = [nbformat.v4.new_output(
        output_type="stream", name="stdout",
        text="Baseline rate (4Q avg): 0.0585\n",
    )]
    nb.cells.append(c3)

    # Cell 4: Stress scenarios (code)
    c4 = nbformat.v4.new_code_cell(
        "# Macro stress scenarios\n"
        "# NOTE: Overlay factors are Alice's expert estimates — not from any formal model\n"
        "#       These have never been peer-reviewed or validated.\n"
        "STRESS_SCENARIOS = {\n"
        "    \"base\":                {\"gdp_shock\": 0.0,  \"unemp_shock\": 0.0, \"overlay\": 0.000},\n"
        "    \"mild_downturn\":       {\"gdp_shock\": -0.5, \"unemp_shock\": 0.3, \"overlay\": 0.015},\n"
        "    \"moderate_recession\":  {\"gdp_shock\": -1.5, \"unemp_shock\": 1.0, \"overlay\": 0.035},\n"
        "    \"severe_recession\":    {\"gdp_shock\": -3.0, \"unemp_shock\": 2.5, \"overlay\": 0.050},\n"
        "}"
    )
    nb.cells.append(c4)

    # Cell 5: Run stress scenarios (code)
    c5 = nbformat.v4.new_code_cell(
        "# Run stress scenarios\n"
        "results = []\n"
        "for name, params in STRESS_SCENARIOS.items():\n"
        "    adjusted_rate = baseline_rate + params[\"overlay\"]\n"
        "    results.append({\n"
        "        \"scenario\": name,\n"
        "        \"baseline\": round(baseline_rate, 4),\n"
        "        \"overlay\": params[\"overlay\"],\n"
        "        \"stressed_rate\": round(adjusted_rate, 4),\n"
        "    })\n\n"
        "stress_df = pd.DataFrame(results)\n"
        "print(stress_df.to_string(index=False))"
    )
    c5.outputs = [nbformat.v4.new_output(
        output_type="stream", name="stdout",
        text=(
            "           scenario  baseline  overlay  stressed_rate\n"
            "               base    0.0585    0.000         0.0585\n"
            "      mild_downturn    0.0585    0.015         0.0735\n"
            "  moderate_recession    0.0585    0.035         0.0935\n"
            "   severe_recession    0.0585    0.050         0.1085\n"
        ),
    )]
    nb.cells.append(c5)

    # Cell 6: Load forecast for comparison (code)
    c6 = nbformat.v4.new_code_cell(
        "# Compare with actual forecast (Alice's manual overlays)\n"
        "forecast_df = pd.read_excel(\"Q3_2024_forecast.xlsx\")\n"
        "print(forecast_df)"
    )
    nb.cells.append(c6)

    # Cell 7: Notes (markdown)
    nb.cells.append(nbformat.v4.new_markdown_cell(
        "## Notes\n\n"
        "- Baseline uses **4-quarter average** for speed — policy says 12-quarter\n"
        "  (see `risk_queries.sql` for both variants)\n"
        "- Overlay factors in `STRESS_SCENARIOS` are Alice's expert judgment estimates\n"
        "- **TODO:** Get overlay factors peer-reviewed before year-end\n"
        "- The severe recession overlay (+5%) matches the overlay cap in `portfolio_risk.db`\n"
        "  (`policy_snapshot.overlay_cap = 0.05`)"
    ))

    nbformat.write(nb, output_dir / "stress_testing.ipynb")


# ---------------------------------------------------------------------------
# 8. run_notes.txt
# ---------------------------------------------------------------------------
def generate_run_notes(output_dir: Path) -> None:
    content = """\
RISK FORECAST RUN NOTES — Q3 2024
===================================
Alice Chen — last updated Sept 20, 2024

SEPT 15 RUN (FINAL):
- Ran loss_forecast_model.py against portfolio_risk.db
- Used 4-quarter lookback (shortcut) since GDP is still > 1%
- Baseline came out at 5.85%, consistent with recent trend
- Applied overlays manually in Q3_2024_forecast.xlsx:
    * Near prime:     +1.2%  (two months of >25% cohort variance)
    * Subprime:       +3.1%  (DQ trend + Jan/Feb cohort miss)
    * Deep subprime:  +4.5%  (keeping new product buffer from March)
- Total reserve impact: $17.6M
- Sent board_risk_presentation.pptx to David for review
- Need CFO sign-off on subprime overlay (>$5M threshold)

ESCALATION NOTES:
- Routine (<$2M): just log it in Q3_2024_forecast.xlsx adjustment tab
- Risk committee ($2-5M): present at weekly sync — bring cohort data
- CFO ($5M+): email with loss_forecast_model.py output + rationale
- Board threshold: STILL not formally set. Operating on verbal ~$10M
  from CFO conversation in June. Need to get this in writing.
- If I'm not available, check with Marcus Park — he knows the model
  but hasn't done overlays or escalation himself yet.

OVERLAY DECISION PROCESS (my mental model):
- Watch 30-day DQ rates monthly by segment
- If DQ jumps >15% MoM for any segment, start considering an overlay
- Cohort variance >25% for 2+ consecutive months = apply overlay
- Sizing: roughly 1% overlay per 10% DQ increase, cap at 5%
- New products < 6 months old: always add 20% buffer on top
- Document in Q3_2024_forecast.xlsx column F (I've been bad about this)
- When to remove: 3 consecutive months of <10% variance

LEGACY VS POLICY METHOD:
- risk_queries.sql has both 12q (policy) and 4q (shortcut) queries
- I use the 4-quarter shortcut when GDP > 1.5% and no major Fed moves
- When macro is volatile, switch to 12-quarter (it's more conservative)
- For the board deck I used the 4q shortcut — faster turnaround
- The difference is usually < 0.3% but can hit 1-2% in a regime change
- model_methodology.docx says 12q is required but I've been using 4q
  for 2+ years with no issues raised

TODO BEFORE LEAVING:
- [ ] Document overlay criteria in model_methodology.docx (Section 4 is blank)
- [ ] Get board notification threshold formalized (still verbal)
- [ ] Train Marcus on the overlay process end-to-end
- [ ] Update stress_testing.ipynb to use 12-quarter baseline
- [ ] Make sure board_risk_presentation.pptx methodology slides are accurate
- [ ] Add approval tracking to Q3_2024_forecast.xlsx (column E always empty)
- [ ] Reconcile loss_forecast_model.py with policy-approved parameters
"""
    (output_dir / "run_notes.txt").write_text(content)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> None:
    parser = argparse.ArgumentParser(description="Generate demo data for Golden Gate")
    parser.add_argument(
        "--output-dir", type=Path, default=DEFAULT_OUTPUT,
        help="Directory to write generated files (default: data/)",
    )
    args = parser.parse_args()

    output_dir: Path = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)

    generators = [
        ("loss_forecast_model.py",       generate_python_model),
        ("Q3_2024_forecast.xlsx",        generate_excel_forecast),
        ("portfolio_risk.db",            generate_sqlite_db),
        ("risk_queries.sql",             generate_sql_queries),
        ("board_risk_presentation.pptx", generate_pptx_deck),
        ("model_methodology.docx",       generate_docx_methodology),
        ("stress_testing.ipynb",         generate_notebook),
        ("run_notes.txt",               generate_run_notes),
    ]

    print(f"Generating {len(generators)} demo files in {output_dir}/\n")
    for name, gen_fn in generators:
        gen_fn(output_dir)
        fpath = output_dir / name
        size = fpath.stat().st_size
        print(f"  [OK] {name:<35s} ({size:>8,d} bytes)")

    print(f"\nDone. {len(generators)} files generated.")


if __name__ == "__main__":
    main()
