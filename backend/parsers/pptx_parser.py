"""PPTX parser using python-pptx library."""

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.parsers import ParseResult, register_parser


@register_parser("pptx")
def parse_pptx(path: str) -> ParseResult:
    """Parse a PowerPoint file and extract text, metadata, and references.

    Args:
        path: Path to the PPTX file

    Returns:
        ParseResult with content, metadata, and references
    """
    prs = Presentation(path)

    # Build metadata
    slides_metadata = []
    content_parts = []
    all_text_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract slide title
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Check for notes
        has_notes = False
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            has_notes = bool(notes_text)

        # Count shapes
        shape_count = len(slide.shapes)

        slides_metadata.append({
            "title": title,
            "has_notes": has_notes,
            "shape_count": shape_count
        })

        # Build content for this slide
        slide_title = title if title else "Untitled"
        content_parts.append(f"## Slide {slide_num}: {slide_title}\n")

        # Process shapes
        slide_text = _process_shapes(slide.shapes)
        if slide_text:
            content_parts.append(slide_text)
            all_text_content.append(slide_text)

        # Add speaker notes as blockquotes
        if has_notes:
            blockquote_lines = [f"> {line}" for line in notes_text.split("\n")]
            notes_section = "\n".join(blockquote_lines)
            content_parts.append(f"\n{notes_section}\n")
            all_text_content.append(notes_text)

        content_parts.append("\n")

    # Combine all content
    content = "\n".join(content_parts).strip()

    # Extract references from all text content
    all_text = "\n".join(all_text_content)
    reference_pattern = re.compile(
        r"\b[\w.-]+\.(xlsx|xls|csv|py|sql|ipynb|pdf|pptx|docx|txt|db)\b",
        re.IGNORECASE,
    )
    seen = set()
    references = []
    for m in reference_pattern.finditer(all_text):
        ref = m.group(0)
        if ref.lower() not in seen:
            seen.add(ref.lower())
            references.append(ref)

    # Build metadata dict
    metadata = {
        "slide_count": len(prs.slides),
        "slides": slides_metadata
    }

    # Build result
    file_path = Path(path).resolve()
    result = ParseResult(
        file_id=file_path.name,
        file_path=str(file_path),
        file_type="pptx",
        content=content,
        metadata=metadata,
        references=references
    )

    return result


def _process_shapes(shapes) -> str:
    """Process shapes and extract text content.

    Args:
        shapes: Collection of shapes to process

    Returns:
        Formatted text content from shapes
    """
    parts = []

    for shape in shapes:
        # Skip title shapes (already handled)
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    continue
            except Exception:
                pass

        # Handle group shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if hasattr(shape, "shapes"):
                group_text = _process_shapes(shape.shapes)
                if group_text:
                    parts.append(group_text)
            continue

        # Handle tables
        if shape.has_table:
            table_text = _format_table(shape.table)
            if table_text:
                parts.append(table_text)
            continue

        # Handle text frames
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

    return "\n\n".join(parts)


def _format_table(table) -> str:
    """Format a table as markdown.

    Args:
        table: python-pptx Table object

    Returns:
        Markdown-formatted table string
    """
    if not table.rows:
        return ""

    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
            # Escape pipe characters in cell content
            cell_text = cell_text.replace("|", "\\|")
            cells.append(cell_text)
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Add header separator after first row
    if len(rows) > 1:
        num_cols = len(table.rows[0].cells)
        separator = "| " + " | ".join(["---"] * num_cols) + " |"
        rows.insert(1, separator)

    return "\n".join(rows)
